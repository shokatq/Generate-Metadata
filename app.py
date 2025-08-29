import os
import json
import logging
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import asyncio
from io import BytesIO
import threading
from functools import wraps
import time
import re
from collections import defaultdict
import uuid

# File processing libraries
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import pandas as pd

# Azure libraries
from azure.storage.blob import BlobServiceClient
from azure.cosmos import CosmosClient
from azure.cosmos.exceptions import CosmosResourceNotFoundError
from openai import AzureOpenAI

# Flask
from flask import Flask, request, jsonify
from werkzeug.exceptions import BadRequest, NotFound, InternalServerError

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class FileMetadataGenerator:
    def __init__(self):
        # Azure Blob Storage
        self.blob_service_client = BlobServiceClient.from_connection_string(
            os.getenv('AZURE_STORAGE_CONNECTION_STRING_1')
        )
        self.container_name = "weezyaifiles"
        
        # Azure Cosmos DB
        self.cosmos_client = CosmosClient(
            url=os.getenv('COSMOS_ENDPOINT'),
            credential=os.getenv('COSMOS_KEY')
        )
        self.database = self.cosmos_client.get_database_client('weezyai')
        self.container = self.database.get_container_client('files')
        
        # Azure OpenAI for embeddings
        self.openai_embedding_client = AzureOpenAI(
            api_key=os.getenv('OPENAI_API_KEY'),
            api_version="2024-12-01-preview",
            azure_endpoint="https://weez-openai-resource.openai.azure.com/"
        )
        
        # Azure OpenAI for text generation (GPT-4o)
        self.openai_text_client = AzureOpenAI(
            api_key=os.getenv('OPENAI_API_KEY'),
            api_version="2024-12-01-preview",
            azure_endpoint="https://weez-openai-resource.openai.azure.com/"
        )
        
        # Supported file types
        self.supported_extensions = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.xls': 'application/vnd.ms-excel',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.txt': 'text/plain'
        }
        
        # Fields that indicate new metadata format
        self.new_metadata_fields = ['textSummary', 'embedding', 'processed_at']
        
        # Chunk size for embeddings (optimize for text-embedding-3-large)
        self.chunk_size = 1000  # Characters per chunk
        self.chunk_overlap = 200  # Overlap between chunks

    def is_metadata_updated(self, metadata: Dict[str, Any]) -> bool:
        """Check if metadata has been updated with new fields"""
        # Check if any of the new metadata fields are present and not empty
        for field in self.new_metadata_fields:
            if field in metadata and metadata[field]:
                return True
        return False

    def get_all_old_metadata(self, user_id: Optional[str] = None) -> List[Dict[str, Any]]:
        """Get all metadata records that need updating"""
        try:
            if user_id:
                # Query for specific user's old metadata
                query = "SELECT * FROM c WHERE c.user_id = @user_id"
                parameters = [{"name": "@user_id", "value": user_id}]
            else:
                # Query for all old metadata
                query = "SELECT * FROM c"
                parameters = []
            
            items = list(self.container.query_items(
                query=query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            # Filter out items that already have updated metadata
            old_metadata_items = []
            for item in items:
                if not self.is_metadata_updated(item):
                    old_metadata_items.append(item)
            
            logger.info(f"Found {len(old_metadata_items)} items that need metadata updates")
            return old_metadata_items
        
        except Exception as e:
            logger.error(f"Error retrieving old metadata: {str(e)}")
            return []

    def download_file_from_blob(self, file_path: str) -> bytes:
        """Download file from Azure Blob Storage with better error handling"""
        try:
            logger.info(f"Attempting to download file from blob: {file_path}")
            
            # Clean the file path - remove any leading slashes
            cleaned_path = file_path.lstrip('/')
            logger.info(f"Cleaned file path: {cleaned_path}")
            
            blob_client = self.blob_service_client.get_blob_client(
                container=self.container_name,
                blob=cleaned_path
            )
            
            # Check if blob exists first
            if not blob_client.exists():
                raise Exception(f"Blob does not exist: {cleaned_path}")
            
            # Get blob properties for debugging
            blob_properties = blob_client.get_blob_properties()
            logger.info(f"Blob properties - Size: {blob_properties.size}, Content-Type: {blob_properties.content_settings.content_type}")
            
            blob_data = blob_client.download_blob()
            content = blob_data.readall()
            
            logger.info(f"Successfully downloaded {len(content)} bytes from blob")
            return content
        
        except Exception as e:
            logger.error(f"Error downloading file from blob: {str(e)}")
            logger.error(f"Container: {self.container_name}, Blob: {file_path}")
            raise Exception(f"Failed to download file from blob storage: {str(e)}")

    def clean_text(self, text: str) -> str:
        """Clean and normalize text for better processing"""
        # Remove excessive whitespace and normalize
        text = re.sub(r'\s+', ' ', text.strip())
        
        # Remove special characters that might interfere with processing
        text = re.sub(r'[^\w\s\.,!?;:()\-\'"\/]', ' ', text)
        
        # Normalize quotes
        text = re.sub(r'["""]', '"', text)
        text = re.sub(r"[''']", "'", text)

        return text

    def create_text_chunks(self, text: str, metadata: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """Create overlapping text chunks for better embedding accuracy"""
        chunks = []
        text = self.clean_text(text)
        
        # If text is short, return as single chunk
        if len(text) <= self.chunk_size:
            chunks.append({
                'text': text,
                'chunk_index': 0,
                'start_char': 0,
                'end_char': len(text),
                'metadata': metadata or {}
            })
            return chunks
        
        # Create overlapping chunks
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            
            # Try to break at sentence boundaries
            if end < len(text):
                # Look for sentence endings within overlap range
                sentence_end = max(
                    text.rfind('.', start, end),
                    text.rfind('!', start, end),
                    text.rfind('?', start, end)
                )
                
                if sentence_end > start:
                    end = sentence_end + 1
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunks.append({
                    'text': chunk_text,
                    'chunk_index': chunk_index,
                    'start_char': start,
                    'end_char': end,
                    'metadata': metadata or {}
                })
                chunk_index += 1
            
            # Move start position with overlap
            start = max(start + self.chunk_size - self.chunk_overlap, end)
            
            if start >= len(text):
                break
        
        logger.info(f"Created {len(chunks)} text chunks")
        return chunks

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file with improved structure preservation"""
        try:
            pdf_file = BytesIO(file_content)
            reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                # Add page break indicators
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page_text + "\n"
            
            return self.clean_text(text)
        
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return ""

    def extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file with structure preservation"""
        try:
            doc_file = BytesIO(file_content)
            doc = docx.Document(doc_file)
            text = ""
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Preserve paragraph structure
                    text += para.text + "\n\n"
            
            # Extract table content
            for table in doc.tables:
                text += "\n--- Table ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    text += row_text + "\n"
                text += "--- End Table ---\n\n"
            
            return self.clean_text(text)
        
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return ""

    def extract_text_from_xlsx(self, file_content: bytes) -> str:
        """Extract text from XLSX file with improved structure"""
        try:
            excel_file = BytesIO(file_content)
            workbook = openpyxl.load_workbook(excel_file, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Get headers from first row
                headers = []
                first_row = next(sheet.iter_rows(min_row=1, max_row=1, values_only=True), None)
                if first_row:
                    headers = [str(cell) if cell is not None else "" for cell in first_row]
                    text += "Headers: " + " | ".join(headers) + "\n"
                
                # Extract data rows
                for row_num, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), 2):
                    if any(cell is not None for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        text += f"Row {row_num}: {row_text}\n"
                
                text += f"--- End Sheet: {sheet_name} ---\n\n"
            
            return self.clean_text(text)
        
        except Exception as e:
            logger.error(f"Error extracting text from XLSX: {str(e)}")
            return ""

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX file with slide structure"""
        try:
            ppt_file = BytesIO(file_content)
            presentation = Presentation(ppt_file)
            text = ""
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                
                text += f"--- End Slide {slide_num} ---\n\n"
            
            return self.clean_text(text)
        
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            return ""

    def extract_text_from_txt(self, file_content: bytes) -> str:
        """Extract text from TXT file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    logger.info(f"Successfully decoded text file using {encoding} encoding")
                    return self.clean_text(text)
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use utf-8 with errors='ignore'
            text = file_content.decode('utf-8', errors='ignore')
            logger.warning("Used utf-8 with errors='ignore' for text file")
            return self.clean_text(text)
        
        except Exception as e:
            logger.error(f"Error extracting text from TXT: {str(e)}")
            return ""

    def extract_text_from_file(self, file_content: bytes, file_extension: str) -> str:
        """Extract text based on file extension"""
        extraction_methods = {
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.doc': self.extract_text_from_docx,
            '.xlsx': self.extract_text_from_xlsx,
            '.xls': self.extract_text_from_xlsx,
            '.pptx': self.extract_text_from_pptx,
            '.ppt': self.extract_text_from_pptx,
            '.txt': self.extract_text_from_txt
        }
        
        method = extraction_methods.get(file_extension.lower())
        if method:
            return method(file_content)
        else:
            logger.warning(f"Unsupported file extension: {file_extension}")
            return ""

    def generate_chunk_embeddings(self, text_chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Generate embeddings for text chunks"""
        try:
            chunk_embeddings = []
            
            for chunk in text_chunks:
                # Generate embedding for chunk
                response = self.openai_embedding_client.embeddings.create(
                    model="text-embedding-3-large",
                    input=chunk['text'],
                    encoding_format="float"
                )
                
                embedding_data = {
                    'chunk_index': chunk['chunk_index'],
                    'text': chunk['text'],
                    'embedding': response.data[0].embedding,
                    'start_char': chunk['start_char'],
                    'end_char': chunk['end_char'],
                    'text_length': len(chunk['text']),
                    'metadata': chunk['metadata']
                }
                
                chunk_embeddings.append(embedding_data)
            
            logger.info(f"Generated embeddings for {len(chunk_embeddings)} chunks")
            return chunk_embeddings
        
        except Exception as e:
            logger.error(f"Error generating chunk embeddings: {str(e)}")
            return []

    def store_chunk_documents(self, file_id: str, file_name: str, user_id: str, 
                             chunk_embeddings: List[Dict[str, Any]], 
                             base_metadata: Dict[str, Any]) -> int:
        """Store each chunk as a separate document in Cosmos DB"""
        try:
            stored_count = 0
            current_time = datetime.utcnow().isoformat() + 'Z'
            
            # Extract base metadata fields
            created_at = base_metadata.get('created_at', current_time)
            platform = base_metadata.get('platform', 'platform_sync')
            mime_type = base_metadata.get('mime_type', 'unknown')
            
            for chunk_data in chunk_embeddings:
                # Create unique chunk ID
                chunk_id = f"{file_id}-chunk-{chunk_data['chunk_index']}"
                
                # Create chunk document
                chunk_document = {
                    'id': chunk_id,
                    'chunk_index': chunk_data['chunk_index'],
                    'text': chunk_data['text'],
                    'embedding': chunk_data['embedding'],
                    'file_id': file_id,
                    'fileName': file_name,
                    'user_id': user_id,
                    'platform': platform,
                    'created_at': created_at,
                    'mime_type': mime_type,
                    'processed_at': current_time,
                    'source': 'platform_sync',
                    'summary_type': None,
                    'metadata': {
                        'start_char': chunk_data['start_char'],
                        'end_char': chunk_data['end_char'],
                        'text_length': chunk_data['text_length'],
                        'embedding_model': 'text-embedding-3-large',
                        'processing_version': '2.0'
                    }
                }
                
                # Store chunk document
                try:
                    self.container.upsert_item(chunk_document)
                    stored_count += 1
                    logger.info(f"Stored chunk {chunk_data['chunk_index']} for file {file_id}")
                    
                except Exception as e:
                    logger.error(f"Error storing chunk {chunk_data['chunk_index']} for file {file_id}: {str(e)}")
                    continue
            
            logger.info(f"Successfully stored {stored_count} chunks for file {file_id}")
            return stored_count
            
        except Exception as e:
            logger.error(f"Error storing chunk documents: {str(e)}")
            return 0

    def get_existing_metadata(self, file_id: str, user_id: str) -> Dict[str, Any]:
        """Retrieve existing metadata from Cosmos DB"""
        try:
            response = self.container.read_item(
                item=file_id,
                partition_key=user_id
            )
            return response
        except CosmosResourceNotFoundError:
            logger.error(f"Metadata not found for file ID: {file_id} with user ID: {user_id}")
            raise Exception(f"File metadata not found for ID: {file_id}")
        except Exception as e:
            logger.error(f"Error retrieving existing metadata: {str(e)}")
            raise Exception(f"Failed to retrieve metadata: {str(e)}")

    def process_single_file_metadata(self, metadata: Dict[str, Any]) -> Tuple[bool, str, Dict[str, Any]]:
        """Process a single file's metadata and store chunks separately"""
        try:
            file_id = metadata.get('id')
            file_path = metadata.get('filePath')
            filename = metadata.get('fileName', '')
            user_id = metadata.get('user_id')
            
            logger.info(f"Processing file metadata for: {file_id}")
            
            # Check if file exists in blob storage
            try:
                file_content = self.download_file_from_blob(file_path)
                logger.info(f"Successfully downloaded file, size: {len(file_content)} bytes")
            except Exception as e:
                logger.warning(f"Could not download file {file_path}: {str(e)}")
                return False, f"File not accessible: {str(e)}", metadata
            
            # Determine file extension
            file_extension = os.path.splitext(filename)[-1].lower()
            
            if file_extension not in self.supported_extensions:
                logger.warning(f"Unsupported file type: {file_extension}")
                return False, f"Unsupported file type: {file_extension}", metadata
            
            # Extract text from file
            logger.info(f"Extracting text from file: {filename}")
            extracted_text = self.extract_text_from_file(file_content, file_extension)
            
            if not extracted_text.strip():
                logger.warning(f"No text extracted from file: {filename}")
                extracted_text = f"Content from {filename}"
            
            logger.info(f"Extracted text length: {len(extracted_text)} characters")
            
            # Create text chunks
            logger.info("Creating text chunks...")
            text_chunks = self.create_text_chunks(extracted_text, {
                'filename': filename,
                'file_id': file_id,
                'file_type': file_extension
            })
            
            # Generate chunk embeddings
            logger.info("Generating chunk embeddings...")
            chunk_embeddings = self.generate_chunk_embeddings(text_chunks)
            
            if not chunk_embeddings:
                logger.error("No embeddings generated for chunks")
                return False, "Failed to generate embeddings", metadata
            
            # Store chunks as separate documents
            logger.info("Storing chunks as separate documents...")
            stored_count = self.store_chunk_documents(
                file_id=file_id,
                file_name=filename,
                user_id=user_id,
                chunk_embeddings=chunk_embeddings,
                base_metadata=metadata
            )
            
            if stored_count == 0:
                logger.error("No chunks were stored successfully")
                return False, "Failed to store chunks", metadata
            
            # Update original metadata to mark as processed
            updated_metadata = metadata.copy()
            updated_metadata.update({
                'processed_at': datetime.utcnow().isoformat() + 'Z',
                'total_chunks': stored_count,
                'text_length': len(extracted_text),
                'processing_version': '2.0',
                'chunk_storage_format': 'separate_documents'
            })
            
            # Update the original metadata document
            self.container.upsert_item(updated_metadata)
            
            logger.info(f"Successfully processed file {file_id} with {stored_count} chunks")
            return True, "Success", updated_metadata
        
        except Exception as e:
            logger.error(f"Error processing file metadata: {str(e)}")
            return False, str(e), metadata

    def process_file_metadata(self, file_id: str, user_id: str, file_path: str) -> Dict[str, Any]:
        """Process file metadata - main entry point"""
        try:
            logger.info(f"Starting metadata processing for file_id: {file_id}")
            
            # Get existing metadata with user_id as partition key
            existing_metadata = self.get_existing_metadata(file_id, user_id)
            logger.info(f"Retrieved existing metadata for file: {file_id}")
            
            # Check if metadata is already updated
            if self.is_metadata_updated(existing_metadata):
                logger.info(f"Metadata already updated for file: {file_id}")
                return existing_metadata
            
            # Process the file
            success, message, updated_metadata = self.process_single_file_metadata(existing_metadata)
            
            if success:
                logger.info(f"Successfully processed metadata for file: {file_id}")
                return updated_metadata
            else:
                logger.error(f"Failed to process metadata for file: {file_id}, Error: {message}")
                raise Exception(f"Processing failed: {message}")
        
        except Exception as e:
            logger.error(f"Error in process_file_metadata: {str(e)}")
            raise Exception(f"Failed to process file metadata: {str(e)}")

    def bulk_update_metadata(self, user_id: Optional[str] = None, batch_size: int = 10) -> Dict[str, Any]:
        """Bulk update metadata for old records"""
        try:
            logger.info(f"Starting bulk metadata update for user: {user_id or 'all users'}")
            
            # Get all old metadata records
            old_metadata_items = self.get_all_old_metadata(user_id)
            
            if not old_metadata_items:
                logger.info("No items found that need metadata updates")
                return {
                    'total_items': 0,
                    'processed': 0,
                    'successful': 0,
                    'failed': 0,
                    'errors': []
                }
            
            # Process in batches
            total_items = len(old_metadata_items)
            processed = 0
            successful = 0
            failed = 0
            errors = []
            
            logger.info(f"Processing {total_items} items in batches of {batch_size}")
            
            for i in range(0, total_items, batch_size):
                batch = old_metadata_items[i:i + batch_size]
                logger.info(f"Processing batch {i//batch_size + 1}/{(total_items + batch_size - 1)//batch_size}")
                
                for metadata in batch:
                    try:
                        file_id = metadata.get('id')
                        logger.info(f"Processing file: {file_id}")
                        
                        success, message, updated_metadata = self.process_single_file_metadata(metadata)
                        processed += 1
                        
                        if success:
                            successful += 1
                            logger.info(f"Successfully updated: {file_id}")
                        else:
                            failed += 1
                            error_msg = f"Failed to update {file_id}: {message}"
                            logger.error(error_msg)
                            errors.append(error_msg)
                    
                    except Exception as e:
                        processed += 1
                        failed += 1
                        error_msg = f"Exception processing {metadata.get('id', 'unknown')}: {str(e)}"
                        logger.error(error_msg)
                        errors.append(error_msg)
                
                # Add small delay between batches to avoid overwhelming services
                if i + batch_size < total_items:
                    time.sleep(1)
            
            result = {
                'total_items': total_items,
                'processed': processed,
                'successful': successful,
                'failed': failed,
                'errors': errors[:50]  # Limit errors to first 50
            }
            
            logger.info(f"Bulk update completed: {result}")
            return result
        
        except Exception as e:
            logger.error(f"Error in bulk metadata update: {str(e)}")
            raise Exception(f"Bulk update failed: {str(e)}")

    def search_similar_content(self, query_text: str, user_id: str, limit: int = 10, similarity_threshold: float = 0.7) -> List[Dict[str, Any]]:
        """Search for similar content using chunk embeddings"""
        try:
            # Generate embedding for query
            response = self.openai_embedding_client.embeddings.create(
                model="text-embedding-3-large",
                input=query_text,
                encoding_format="float"
            )
            query_embedding = response.data[0].embedding
            
            # Query Cosmos DB for chunk documents
            query = """
                SELECT c.id, c.file_id, c.fileName, c.chunk_index, c.text, c.embedding, c.metadata
                FROM c 
                WHERE c.user_id = @user_id 
                AND IS_DEFINED(c.embedding)
                AND IS_DEFINED(c.chunk_index)
            """
            
            parameters = [{"name": "@user_id", "value": user_id}]
            
            items = list(self.container.query_items(
                query=query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            # Calculate similarities
            results = []
            for item in items:
                chunk_embedding = item.get('embedding', [])
                if chunk_embedding:
                    similarity = self.calculate_cosine_similarity(query_embedding, chunk_embedding)
                    
                    if similarity >= similarity_threshold:
                        results.append({
                            'chunk_id': item['id'],
                            'file_id': item['file_id'],
                            'filename': item.get('fileName', ''),
                            'chunk_index': item.get('chunk_index', 0),
                            'text': item.get('text', '')[:200] + '...',  # Truncate for response
                            'similarity': similarity,
                            'metadata': item.get('metadata', {})
                        })
            
            # Sort by similarity
            results.sort(key=lambda x: x['similarity'], reverse=True)
            
            return results[:limit]
        
        except Exception as e:
            logger.error(f"Error in similarity search: {str(e)}")
            return []

    def calculate_cosine_similarity(self, vec1: List[float], vec2: List[float]) -> float:
        """Calculate cosine similarity between two vectors"""
        try:
            if len(vec1) != len(vec2):
                return 0.0
            
            dot_product = sum(a * b for a, b in zip(vec1, vec2))
            magnitude1 = sum(a * a for a in vec1) ** 0.5
            magnitude2 = sum(b * b for b in vec2) ** 0.5
            
            if magnitude1 == 0.0 or magnitude2 == 0.0:
                return 0.0
            
            return dot_product / (magnitude1 * magnitude2)
        
        except Exception as e:
            logger.error(f"Error calculating cosine similarity: {str(e)}")
            return 0.0

    def get_file_analytics(self, user_id: str) -> Dict[str, Any]:
        """Get analytics about processed files"""
        try:
            # Get file-level analytics
            file_query = """
                SELECT 
                    COUNT(1) as total_files,
                    SUM(c.text_length) as total_text_length,
                    SUM(c.total_chunks) as total_chunks,
                    AVG(c.text_length) as avg_text_length,
                    AVG(c.total_chunks) as avg_chunks_per_file
                FROM c 
                WHERE c.user_id = @user_id 
                AND IS_DEFINED(c.processed_at)
                AND NOT IS_DEFINED(c.chunk_index)
            """
            
            # Get chunk-level analytics
            chunk_query = """
                SELECT 
                    COUNT(1) as total_chunks_stored
                FROM c 
                WHERE c.user_id = @user_id 
                AND IS_DEFINED(c.chunk_index)
            """
            
            parameters = [{"name": "@user_id", "value": user_id}]
            
            file_results = list(self.container.query_items(
                query=file_query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            chunk_results = list(self.container.query_items(
                query=chunk_query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            # Get file type distribution
            file_type_query = """
                SELECT c.mime_type, COUNT(1) as count
                FROM c 
                WHERE c.user_id = @user_id 
                AND IS_DEFINED(c.processed_at)
                AND NOT IS_DEFINED(c.chunk_index)
                GROUP BY c.mime_type
            """
            
            file_type_results = list(self.container.query_items(
                query=file_type_query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            # Compile analytics
            analytics = {
                'user_id': user_id,
                'file_analytics': file_results[0] if file_results else {},
                'chunk_analytics': chunk_results[0] if chunk_results else {},
                'file_type_distribution': file_type_results,
                'generated_at': datetime.utcnow().isoformat() + 'Z'
            }
            
            return analytics
            
        except Exception as e:
            logger.error(f"Error getting file analytics: {str(e)}")
            return {
                'user_id': user_id,
                'error': str(e),
                'generated_at': datetime.utcnow().isoformat() + 'Z'
            }

    def delete_old_embeddings(self, user_id: str, file_id: str = None) -> Dict[str, Any]:
        """Delete old embedding format data"""
        try:
            if file_id:
                # Delete specific file's old embeddings
                query = """
                    SELECT c.id FROM c 
                    WHERE c.user_id = @user_id 
                    AND c.file_id = @file_id
                    AND IS_DEFINED(c.embedding)
                    AND IS_ARRAY(c.embedding)
                    AND NOT IS_DEFINED(c.chunk_index)
                """
                parameters = [
                    {"name": "@user_id", "value": user_id},
                    {"name": "@file_id", "value": file_id}
                ]
            else:
                # Delete all old embeddings for user
                query = """
                    SELECT c.id FROM c 
                    WHERE c.user_id = @user_id 
                    AND IS_DEFINED(c.embedding)
                    AND IS_ARRAY(c.embedding)
                    AND NOT IS_DEFINED(c.chunk_index)
                """
                parameters = [{"name": "@user_id", "value": user_id}]
            
            items_to_delete = list(self.container.query_items(
                query=query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            deleted_count = 0
            for item in items_to_delete:
                try:
                    self.container.delete_item(
                        item=item['id'],
                        partition_key=user_id
                    )
                    deleted_count += 1
                except Exception as e:
                    logger.error(f"Error deleting item {item['id']}: {str(e)}")
                    continue
            
            return {
                'deleted_count': deleted_count,
                'total_found': len(items_to_delete),
                'user_id': user_id,
                'file_id': file_id
            }
            
        except Exception as e:
            logger.error(f"Error deleting old embeddings: {str(e)}")
            return {
                'error': str(e),
                'deleted_count': 0,
                'user_id': user_id,
                'file_id': file_id
            }

    def get_chunk_content(self, chunk_id: str, user_id: str) -> Dict[str, Any]:
        """Get full chunk content by chunk ID"""
        try:
            response = self.container.read_item(
                item=chunk_id,
                partition_key=user_id
            )
            return response
        except CosmosResourceNotFoundError:
            logger.error(f"Chunk not found: {chunk_id}")
            return None
        except Exception as e:
            logger.error(f"Error retrieving chunk content: {str(e)}")
            return None

    def get_file_chunks(self, file_id: str, user_id: str) -> List[Dict[str, Any]]:
        """Get all chunks for a specific file"""
        try:
            query = """
                SELECT c.id, c.chunk_index, c.text, c.metadata, c.start_char, c.end_char
                FROM c 
                WHERE c.user_id = @user_id 
                AND c.file_id = @file_id
                AND IS_DEFINED(c.chunk_index)
                ORDER BY c.chunk_index
            """
            
            parameters = [
                {"name": "@user_id", "value": user_id},
                {"name": "@file_id", "value": file_id}
            ]
            
            chunks = list(self.container.query_items(
                query=query,
                parameters=parameters,
                enable_cross_partition_query=True
            ))
            
            return chunks
            
        except Exception as e:
            logger.error(f"Error getting file chunks: {str(e)}")
            return []

    def health_check(self) -> Dict[str, Any]:
        """Perform health check on all services"""
        try:
            health_status = {
                'timestamp': datetime.utcnow().isoformat() + 'Z',
                'services': {}
            }
            
            # Check Blob Storage
            try:
                containers = list(self.blob_service_client.list_containers())
                health_status['services']['blob_storage'] = {
                    'status': 'healthy',
                    'containers_count': len(containers)
                }
            except Exception as e:
                health_status['services']['blob_storage'] = {
                    'status': 'unhealthy',
                    'error': str(e)
                }
            
            # Check Cosmos DB
            try:
                container_props = self.container.read()
                health_status['services']['cosmos_db'] = {
                    'status': 'healthy',
                    'container_id': container_props['id']
                }
            except Exception as e:
                health_status['services']['cosmos_db'] = {
                    'status': 'unhealthy',
                    'error': str(e)
                }
            
            # Check OpenAI Embedding Service
            try:
                test_response = self.openai_embedding_client.embeddings.create(
                    model="text-embedding-3-large",
                    input="test",
                    encoding_format="float"
                )
                health_status['services']['openai_embedding'] = {
                    'status': 'healthy',
                    'embedding_dim': len(test_response.data[0].embedding)
                }
            except Exception as e:
                health_status['services']['openai_embedding'] = {
                    'status': 'unhealthy',
                    'error': str(e)
                }
            
            # Check OpenAI Text Service
            try:
                test_response = self.openai_text_client.chat.completions.create(
                    model="gpt-4o",
                    messages=[{"role": "user", "content": "test"}],
                    max_tokens=1
                )
                health_status['services']['openai_text'] = {
                    'status': 'healthy',
                    'model': test_response.model
                }
            except Exception as e:
                health_status['services']['openai_text'] = {
                    'status': 'unhealthy',
                    'error': str(e)
                }
            
            # Overall health
            all_healthy = all(
                service['status'] == 'healthy' 
                for service in health_status['services'].values()
            )
            health_status['overall_status'] = 'healthy' if all_healthy else 'unhealthy'
            
            return health_status
            
        except Exception as e:
            logger.error(f"Error in health check: {str(e)}")
            return {
                'timestamp': datetime.utcnow().isoformat() + 'Z',
                'overall_status': 'unhealthy',
                'error': str(e)
            }


# Flask Application
app = Flask(__name__)
metadata_generator = FileMetadataGenerator()

def handle_errors(f):
    """Decorator to handle errors consistently"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        try:
            return f(*args, **kwargs)
        except Exception as e:
            logger.error(f"Error in {f.__name__}: {str(e)}")
            return jsonify({
                'error': str(e),
                'success': False,
                'timestamp': datetime.utcnow().isoformat() + 'Z'
            }), 500
    return decorated_function

@app.route('/health', methods=['GET'])
@handle_errors
def health_check():
    """Health check endpoint"""
    health_status = metadata_generator.health_check()
    status_code = 200 if health_status['overall_status'] == 'healthy' else 503
    return jsonify(health_status), status_code

@app.route('/process-file', methods=['POST'])
@handle_errors
def process_file():
    """Process a single file's metadata"""
    data = request.get_json()
    
    if not data:
        raise BadRequest("No JSON data provided")
    
    file_id = data.get('file_id')
    user_id = data.get('user_id')
    file_path = data.get('file_path')
    
    if not all([file_id, user_id, file_path]):
        raise BadRequest("Missing required fields: file_id, user_id, file_path")
    
    logger.info(f"Processing file request: {file_id} for user: {user_id}")
    
    result = metadata_generator.process_file_metadata(file_id, user_id, file_path)
    
    return jsonify({
        'success': True,
        'file_id': file_id,
        'result': result,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.route('/bulk-update', methods=['POST'])
@handle_errors
def bulk_update():
    """Bulk update metadata for multiple files"""
    data = request.get_json() or {}
    
    user_id = data.get('user_id')
    batch_size = data.get('batch_size', 10)
    
    if batch_size > 50:
        batch_size = 50  # Limit batch size
    
    logger.info(f"Starting bulk update for user: {user_id or 'all users'}")
    
    result = metadata_generator.bulk_update_metadata(user_id, batch_size)
    
    return jsonify({
        'success': True,
        'result': result,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.route('/search', methods=['POST'])
@handle_errors
def search_content():
    """Search for similar content"""
    data = request.get_json()
    
    if not data:
        raise BadRequest("No JSON data provided")
    
    query_text = data.get('query')
    user_id = data.get('user_id')
    limit = data.get('limit', 10)
    similarity_threshold = data.get('similarity_threshold', 0.7)
    
    if not all([query_text, user_id]):
        raise BadRequest("Missing required fields: query, user_id")
    
    if limit > 50:
        limit = 50  # Limit results
    
    logger.info(f"Searching content for user: {user_id}, query: {query_text[:50]}...")
    
    results = metadata_generator.search_similar_content(
        query_text, user_id, limit, similarity_threshold
    )
    
    return jsonify({
        'success': True,
        'results': results,
        'total_results': len(results),
        'query': query_text,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.route('/analytics', methods=['GET'])
@handle_errors
def get_analytics():
    """Get file analytics for a user"""
    user_id = request.args.get('user_id')
    
    if not user_id:
        raise BadRequest("Missing required parameter: user_id")
    
    logger.info(f"Getting analytics for user: {user_id}")
    
    analytics = metadata_generator.get_file_analytics(user_id)
    
    return jsonify({
        'success': True,
        'analytics': analytics,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.route('/file-chunks', methods=['GET'])
@handle_errors
def get_file_chunks():
    """Get all chunks for a specific file"""
    file_id = request.args.get('file_id')
    user_id = request.args.get('user_id')
    
    if not all([file_id, user_id]):
        raise BadRequest("Missing required parameters: file_id, user_id")
    
    logger.info(f"Getting chunks for file: {file_id}, user: {user_id}")
    
    chunks = metadata_generator.get_file_chunks(file_id, user_id)
    
    return jsonify({
        'success': True,
        'file_id': file_id,
        'chunks': chunks,
        'total_chunks': len(chunks),
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.route('/chunk-content', methods=['GET'])
@handle_errors
def get_chunk_content():
    """Get full content of a specific chunk"""
    chunk_id = request.args.get('chunk_id')
    user_id = request.args.get('user_id')
    
    if not all([chunk_id, user_id]):
        raise BadRequest("Missing required parameters: chunk_id, user_id")
    
    logger.info(f"Getting chunk content: {chunk_id}, user: {user_id}")
    
    chunk = metadata_generator.get_chunk_content(chunk_id, user_id)
    
    if not chunk:
        raise NotFound(f"Chunk not found: {chunk_id}")
    
    return jsonify({
        'success': True,
        'chunk': chunk,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.route('/delete-old-embeddings', methods=['DELETE'])
@handle_errors
def delete_old_embeddings():
    """Delete old embedding format data"""
    data = request.get_json() or {}
    
    user_id = data.get('user_id')
    file_id = data.get('file_id')
    
    if not user_id:
        raise BadRequest("Missing required field: user_id")
    
    logger.info(f"Deleting old embeddings for user: {user_id}, file: {file_id or 'all'}")
    
    result = metadata_generator.delete_old_embeddings(user_id, file_id)
    
    return jsonify({
        'success': True,
        'result': result,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    })

@app.errorhandler(BadRequest)
def handle_bad_request(e):
    return jsonify({
        'error': str(e),
        'success': False,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    }), 400

@app.errorhandler(NotFound)
def handle_not_found(e):
    return jsonify({
        'error': str(e),
        'success': False,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    }), 404

@app.errorhandler(InternalServerError)
def handle_internal_error(e):
    return jsonify({
        'error': 'Internal server error',
        'success': False,
        'timestamp': datetime.utcnow().isoformat() + 'Z'
    }), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
